import win32com.client
from pptx import Presentation
from pptx.util import Inches

def extract_excel_data(excel_file):
    excel = win32com.client.Dispatch("Excel.Application")
    wb = excel.Workbooks.Open(excel_file)
    ws = wb.ActiveSheet

    # Find the last row and column with data
    last_row = ws.Cells(ws.Rows.Count, 1).End(-4162).Row  # -4162 corresponds to xlUp
    last_col = ws.Cells(1, ws.Columns.Count).End(-4159).Column  # -4159 corresponds to xlToLeft

    # Extract data into a list of lists
    data = []
    for row in range(1, last_row + 1):
        row_data = []
        for col in range(1, last_col + 1):
            row_data.append(ws.Cells(row, col).Value)
        data.append(row_data)

    # Close workbook and Excel application
    wb.Close(False)
    excel.Quit()

    return data

def create_ppt_table(data):
    prs = Presentation()

    max_rows = 6
    max_cols = 4
    slide_idx = 0
    while len(data) > 0:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide_title = slide.shapes.title
        slide_title.text = f"Table {slide_idx + 1}"

        num_rows = min(max_rows, len(data))
        num_cols = min(max_cols, len(data[0]))

        slide_data = data[:num_rows]

        left = Inches(1)
        top = Inches(1.5)
        width = Inches(8)
        height = Inches(5)
        table = slide.shapes.add_table(num_rows, num_cols, left, top, width, height).table

        for row_idx, row_data in enumerate(slide_data):
            for col_idx, cell_value in enumerate(row_data):
                table.cell(row_idx, col_idx).text = str(cell_value)

        data = data[num_rows:]

        slide_idx += 1

    prs.save("Task_2 Output.pptx")
    print("PowerPoint presentation created successfully!")

if __name__ == "__main__":
    excel_file = r"C:\\Users\\himan\Downloads\\Python Assessment\\Python Assessment Zip\\PPT Assessment Table.xlsx"
    excel_data = extract_excel_data(excel_file)
    create_ppt_table(excel_data)

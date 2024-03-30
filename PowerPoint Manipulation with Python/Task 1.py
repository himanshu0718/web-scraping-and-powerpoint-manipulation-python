import win32com.client as win32
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create an instance of Excel
excel = win32.Dispatch("Excel.Application")

# Open the workbook
workbook = excel.Workbooks.Open(r"C:\\Users\\himan\Downloads\\Python Assessment\\Python Assessment Zip\\PPT Assessment Resource.xlsx")
sheet = workbook.ActiveSheet

# Get the number of rows in the sheet
num_rows = sheet.UsedRange.Rows.Count

# Create a PowerPoint presentation
prs = Presentation()

# Iterate through each row in the Excel sheet
for i in range(2, num_rows + 1):  # Start from row 2, assuming row 1 contains headers
    # Get data for the current row
    simple_string = sheet.Cells(i, 1).Value
    url = sheet.Cells(i, 2).Value
    companies = sheet.Cells(i, 3).Value.split(", ")

    # Create a new slide for each row
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Task 1a: Simple Strings
    title_shape = slide.shapes.title
    title_shape.text = simple_string

    # Task 1b: URLs
    text_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(100), width=Pt(600), height=Pt(50))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = "URL: "
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0, 0, 255)
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = url
    r.hyperlink.address = url

    # Task 1c: Random Companies
    text_box = slide.shapes.add_textbox(left=Pt(50), top=Pt(200), width=Pt(600), height=Pt(300))
    text_frame = text_box.text_frame
    text_frame.word_wrap = True  # Enable text wrapping 
    p = text_frame.add_paragraph()
    p.text = "List of Companies:"
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT
    for company in companies:
        p = text_frame.add_paragraph()
        p.text = company
        p.font.size = Pt(14)

# Save the PowerPoint presentation
prs.save("Task_1 Output.pptx")

# Close Excel
excel.Quit()



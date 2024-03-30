from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def create_chart_slide(prs, data, max_data_points_per_chart):
    """
    Create a slide with a chart based on the provided data.
    If the data exceeds the capacity of one chart, split it across multiple slides.
    """
    slide_layout = prs.slide_layouts[5]  # Use slide layout that contains a title and content
    slide = prs.slides.add_slide(slide_layout)
    
    # Filter out data for missing years
    filtered_data = [(x, y) for x, y in zip(data['x_values'], data['y_values']) if x in data['x_values']]

    chart_data = CategoryChartData()
    chart_data.categories = [x for x, _ in filtered_data]
    chart_data.add_series(None, [y for _, y in filtered_data])

    x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(5)
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data
    ).chart

    chart.has_legend = False  # Remove legend
    chart.chart_style = 16  # Apply a chart style

    category_axis = chart.category_axis
    value_axis = chart.value_axis

    category_axis.tick_labels.font.bold = True
    category_axis.tick_labels.font.size = Pt(12)
    category_axis.tick_labels.font.name = 'Verdana'

    value_axis.tick_labels.font.bold = True
    value_axis.tick_labels.font.size = Pt(12)
    value_axis.tick_labels.font.name = 'Verdana'

    for series in chart.series:
        for point in series.points:
            point.format.fill.solid()
            point.format.fill.fore_color.rgb = RGBColor(0x1F, 0x4E, 0x79)  # "#1F4E79"
    
    # Set the title placeholder text to "Series 1"
    for shape in slide.shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER:
            if shape.placeholder_format.idx == 0:  # Placeholder for title
                shape.text = "Series 1"  # Set title text
    
    return slide

# Sample data
data = {
    'x_values': ['2018', '2019', '2020', '2021', '2022', '2023', '2024', '2025', '2026', '2027'],
    'y_values': [195, 200, 203, 208, 212, 218, 210, 230, 235, 240]
}

# Create a presentation object
prs = Presentation()

# Determine the maximum number of data points per chart
max_data_points_per_chart = 8

# Split the data into chunks to fit into charts
chunks = [
    data['x_values'][i:i + max_data_points_per_chart]
    for i in range(0, len(data['x_values']), max_data_points_per_chart)
]

# Create slides with charts for each chunk of data
for chunk in chunks:
    create_chart_slide(prs, {'x_values': chunk, 'y_values': data['y_values']}, max_data_points_per_chart)

# Save presentation
prs.save('Task_3 Output.pptx')
